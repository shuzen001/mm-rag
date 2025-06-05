import nltk
import os
import subprocess
import tempfile
import uuid 
from langchain_text_splitters import CharacterTextSplitter
from unstructured.partition.pdf import partition_pdf
from unstructured.partition.pptx import partition_pptx
from unstructured.cleaners.core import clean_extra_whitespace
import pathlib
from pdf2image import convert_from_path
import tempfile
from PIL import Image
from pptx import Presentation
import io
import shutil

# 設定一個基礎圖片路徑
base_figure_path = "figures/"

# 將PDF頁面轉換為圖片的新函數
def convert_pdf_to_page_images(file_path, file_name):
    """
    將PDF文件的每一頁轉換為圖片
    file_path: PDF檔案所在資料夾路徑
    file_name: PDF檔案名稱
    回傳：包含(頁碼, 圖片路徑)的列表
    """
    try:
        # 設定資料庫資料夾路徑
        DATABASE_DIR = "./database"
        output_dir = f"{DATABASE_DIR}/figures/{os.path.splitext(file_name)[0]}/pages"
        
        # 確保輸出資料夾存在
        os.makedirs(output_dir, exist_ok=True)
        
        # 完整的PDF路徑
        full_path = os.path.join(file_path, file_name)
        
        print(f"🔄 開始處理 PDF: {full_path}")
        
        # 檢查文件是否存在
        if not os.path.exists(full_path):
            print(f"❌ 錯誤: PDF 文件不存在: {full_path}")
            return []
            
        # 檢查文件大小
        file_size = os.path.getsize(full_path) / (1024 * 1024)  # 轉換為 MB
        print(f"📄 PDF 文件大小: {file_size:.2f} MB")
            
        # 轉換PDF頁面為圖片 (dpi值可以調整，影響圖片質量和大小)
        try:
            print(f"🔄 正在將 {file_name} 轉換為頁面圖片...")
            images = convert_from_path(full_path, dpi=200, thread_count=20)
            print(f"✅ 成功讀取 {len(images)} 頁")
        except Exception as e:
            print(f"❌ 轉換 PDF 頁面時出錯: {str(e)}")
            # 嘗試打印更詳細的錯誤信息
            import traceback
            traceback.print_exc()
            return []
        
        # 儲存每一頁為圖片
        page_image_paths = []
        for i, image in enumerate(images):
            try:
                page_num = i + 1  # 頁碼從1開始
                image_path = os.path.join(output_dir, f'page_{page_num}.jpg')
                image.save(image_path, 'JPEG')
                page_image_paths.append((page_num, image_path))
            except Exception as e:
                print(f"❌ 儲存頁面 {i+1} 時出錯: {str(e)}")
        
        print(f"✅ 已將 {file_name} 的 {len(images)} 頁轉換為圖片")
        return page_image_paths
    except Exception as e:
        print(f"❌ 處理 PDF 文件 {file_name} 時發生未預期的錯誤: {str(e)}")
        # 嘗試打印更詳細的錯誤信息
        import traceback
        traceback.print_exc()
        return []

# 新增幻燈片到圖像的轉換函數
def convert_pptx_to_slide_images(file_path, file_name):
    """
    將 PPTX 文件的每一張幻燈片轉換為圖片
    file_path: PPTX 檔案所在資料夾路徑
    file_name: PPTX 檔案名稱
    回傳：包含(幻燈片編號, 圖片路徑)的列表
    """
    try:
        # 設定資料庫資料夾路徑
        DATABASE_DIR = "./database"
        output_dir = f"{DATABASE_DIR}/figures/{os.path.splitext(file_name)[0]}/slides"
        
        # 確保輸出資料夾存在
        os.makedirs(output_dir, exist_ok=True)
        
        # 完整的 PPTX 路徑
        full_path = os.path.join(file_path, file_name)
        
        print(f"🔄 開始處理 PPTX: {full_path}")
        
        # 檢查文件是否存在
        if not os.path.exists(full_path):
            print(f"❌ 錯誤: PPTX 文件不存在: {full_path}")
            return []
            
        # 檢查文件大小
        file_size = os.path.getsize(full_path) / (1024 * 1024)  # 轉換為 MB
        print(f"📄 PPTX 文件大小: {file_size:.2f} MB")
        
        # 載入 PPTX 文件
        try:
            print(f"🔄 正在載入 {file_name}...")
            presentation = Presentation(full_path)
            print(f"✅ 成功載入 {len(presentation.slides)} 張幻燈片")
        except Exception as e:
            print(f"❌ 載入 PPTX 文件時出錯: {str(e)}")
            import traceback
            traceback.print_exc()
            return []
        
        # 存儲每一張幻燈片為圖片
        slide_image_paths = []
        
        # 由於 python-pptx 不能直接將幻燈片轉為圖像，我們使用臨時檔案處理
        # 這裡我們使用外部命令來轉換 (需要安裝 libreoffice 或類似工具)
        # 或使用替代方案：儲存幻燈片內容的結構化信息
        
        # 轉換方法：將每張幻燈片的內容轉為 HTML 然後渲染
        # 注意：這是一個簡化的實現，實際效果可能不如直接從 PDF 獲取圖像精確
        for i, slide in enumerate(presentation.slides):
            try:
                slide_num = i + 1  # 幻燈片編號從1開始
                slide_path = os.path.join(output_dir, f'slide_{slide_num}.jpg')
                
                # 創建一個空白圖像，設定適當的尺寸
                slide_width, slide_height = 1920, 1080  # 標準 16:9 尺寸
                slide_image = Image.new('RGB', (slide_width, slide_height), color='white')
                
                # 提取幻燈片內容並在圖像上繪製
                # 注意：這只是一個簡化版本，實際上需要更複雜的渲染
                # 為了簡單起見，我們將幻燈片的文本內容保存到圖像中
                from PIL import ImageDraw, ImageFont
                draw = ImageDraw.Draw(slide_image)
                
                # 嘗試獲取一個可用的字體
                try:
                    font = ImageFont.truetype('Arial', 20)
                except IOError:
                    font = ImageFont.load_default()
                
                # 提取幻燈片標題和內容
                content_text = f"Slide {slide_num}"
                y_position = 50
                
                if slide.shapes.title:
                    title_text = slide.shapes.title.text
                    draw.text((50, y_position), f"Title: {title_text}", fill='black', font=font)
                    y_position += 40
                
                # 提取文本框內容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # 繪製每行文本
                        lines = shape.text.split('\n')
                        for line in lines:
                            if len(line.strip()) > 0:
                                draw.text((50, y_position), line, fill='black', font=font)
                                y_position += 30
                
                # 保存圖像
                slide_image.save(slide_path, 'JPEG')
                slide_image_paths.append((slide_num, slide_path))
                print(f"  ✓ 已儲存幻燈片 {slide_num} 圖像")
            
            except Exception as e:
                print(f"❌ 處理幻燈片 {i+1} 時出錯: {str(e)}")
        
        print(f"✅ 已將 {file_name} 的 {len(slide_image_paths)}/{len(presentation.slides)} 張幻燈片轉換為圖片")
        return slide_image_paths
        
    except Exception as e:
        print(f"❌ 處理 PPTX 文件 {file_name} 時發生未預期的錯誤: {str(e)}")
        import traceback
        traceback.print_exc()
        return []

def convert_pptx_to_pdf(file_path, file_name):
    """
    將 PPTX 文件轉換為 PDF 文件，解決中文顯示和檔名問題
    file_path: PPTX 檔案所在資料夾路徑
    file_name: PPTX 檔案名稱
    回傳：PDF 文件的路徑
    """
    try:
        # 完整的 PPTX 路徑
        full_path = os.path.join(file_path, file_name)
        
        # 檢查檔案是否存在
        if not os.path.exists(full_path):
            print(f"❌ 錯誤: PPTX 文件不存在: {full_path}")
            return None
        
        # 檢查檔案權限
        if not os.access(full_path, os.R_OK):
            print(f"❌ 錯誤: 沒有讀取 PPTX 文件的權限: {full_path}")
            try:
                # 嘗試修改權限
                os.chmod(full_path, 0o644)
                print(f"  ✓ 已嘗試修改檔案權限")
            except Exception as e:
                print(f"  ⚠️ 無法修改檔案權限: {e}")
                return None
        
        # 設定臨時目錄 (用於存放轉換後的 PDF)
        temp_dir = tempfile.mkdtemp()
        
        # 建立一個臨時目錄用於存放檔案的 ASCII 版本
        ascii_temp_dir = tempfile.mkdtemp()
        
        # 創建一個檔案名稱的 ASCII 版本 (使用 UUID)
        ascii_file_name = f"pptx_{uuid.uuid4().hex}.pptx"
        ascii_file_path = os.path.join(ascii_temp_dir, ascii_file_name)
        
        # 複製原始檔案到 ASCII 檔名的版本
        try:
            shutil.copy2(full_path, ascii_file_path)
            print(f"  ✓ 已創建臨時檔案副本: {ascii_file_path}")
        except Exception as e:
            print(f"  ❌ 無法創建檔案副本: {e}")
            shutil.rmtree(temp_dir, ignore_errors=True)
            shutil.rmtree(ascii_temp_dir, ignore_errors=True)
            return None
        
        # 生成輸出 PDF 的路徑 (臨時文件)
        pdf_name = f"{os.path.splitext(file_name)[0]}.pdf"
        output_pdf = os.path.join(temp_dir, pdf_name)
        ascii_output_pdf = os.path.join(temp_dir, f"{os.path.splitext(ascii_file_name)[0]}.pdf")
        
        print(f"🔄 開始將 PPTX 轉換為 PDF: {ascii_file_path}")
        
        # 確保系統中有中文字體
        try:
            # 檢查常見中文字體是否存在
            fonts_dir = "/usr/share/fonts"
            found_chinese_font = False
            
            if os.path.exists(fonts_dir):
                for root, dirs, files in os.walk(fonts_dir):
                    for file in files:
                        if file.endswith('.ttf') or file.endswith('.ttc'):
                            # 檢查是否是中文字體
                            if any(keyword in file.lower() for keyword in ['chinese', 'cjk', 'noto', 'wqy', 'droid', 'ming', 'kai', 'heiti', 'song']):
                                print(f"  ✓ 找到中文字體: {os.path.join(root, file)}")
                                found_chinese_font = True
                                break
                    if found_chinese_font:
                        break
            
            if not found_chinese_font:
                print("  ⚠️ 警告：可能找不到合適的中文字體，可能會影響中文顯示")
        except Exception as e:
            print(f"  ⚠️ 檢查中文字體時出錯: {e}")
        
        # 使用 LibreOffice 進行轉換
        try:
            # 設置環境變數，確保 LibreOffice 使用正確的配置
            env = os.environ.copy()
            env['LC_ALL'] = 'C'
            env['LANG'] = 'C'
            
            # 嘗試多種轉換方式
            conversion_methods = [
                # 方法1: 標準轉換
                {
                    'cmd': [
                        'libreoffice', 
                        '--headless', 
                        '--convert-to', 'pdf',
                        '--outdir', temp_dir, 
                        ascii_file_path
                    ],
                    'env': env,
                    'name': "標準方法"
                },
                # 方法2: 使用特定過濾器
                {
                    'cmd': [
                        'libreoffice',
                        '--headless',
                        '--convert-to', 'pdf:writer_pdf_Export',
                        '--outdir', temp_dir,
                        ascii_file_path
                    ],
                    'env': env,
                    'name': "PDF導出過濾器方法"
                },
                # 方法3: 使用PowerPoint專用過濾器
                {
                    'cmd': [
                        'libreoffice',
                        '--headless',
                        '--infilter=impress_MS_PowerPoint_2007_XML',
                        '--convert-to', 'pdf',
                        '--outdir', temp_dir,
                        ascii_file_path
                    ],
                    'env': env,
                    'name': "PowerPoint過濾器方法"
                }
            ]
            
            success = False
            
            # 逐一嘗試不同的轉換方法
            for method in conversion_methods:
                if success:
                    break
                    
                try:
                    print(f"  🔄 嘗試使用{method['name']}轉換: {' '.join(method['cmd'])}")
                    result = subprocess.run(
                        method['cmd'], 
                        capture_output=True, 
                        text=True, 
                        check=False,  # 不要因為命令失敗而拋出異常
                        env=method['env'],
                        timeout=120  # 設置超時，避免卡死
                    )
                    
                    # 查看是否生成了PDF (檢查原始檔名和ASCII檔名的版本)
                    if os.path.exists(ascii_output_pdf):
                        # 如果生成了ASCII檔名的PDF，重命名為原始檔名
                        try:
                            shutil.move(ascii_output_pdf, output_pdf)
                            print(f"  ✓ {method['name']}成功: 已將 {ascii_output_pdf} 重命名為 {output_pdf}")
                            success = True
                            break
                        except Exception as e:
                            print(f"  ⚠️ 重命名PDF文件時出錯: {e}")
                    elif os.path.exists(output_pdf):
                        print(f"  ✓ {method['name']}成功: 生成了 {output_pdf}")
                        success = True
                        break
                    else:
                        print(f"  ❌ {method['name']}失敗: 沒有生成PDF文件")
                        if result.returncode != 0:
                            print(f"  命令返回碼: {result.returncode}")
                        if result.stdout:
                            print(f"  命令輸出: {result.stdout}")
                        if result.stderr:
                            print(f"  命令錯誤: {result.stderr}")
                except subprocess.TimeoutExpired:
                    print(f"  ⚠️ {method['name']}轉換超時")
                except Exception as e:
                    print(f"  ❌ 執行{method['name']}時出錯: {e}")
        
            # 檢查最終結果
            if success and os.path.exists(output_pdf):
                # 檢查 PDF 文件大小
                pdf_size = os.path.getsize(output_pdf) / 1024  # 轉換為 KB
                if pdf_size < 5:  # 小於 5KB 可能是空文件或轉換失敗
                    print(f"  ⚠️ 警告：生成的 PDF 文件非常小 ({pdf_size:.2f} KB)，可能轉換不完整")
                
                # 清理臨時文件
                try:
                    shutil.rmtree(ascii_temp_dir, ignore_errors=True)
                except:
                    pass
                    
                return output_pdf
            else:
                print(f"  ❌ 所有轉換方法都失敗了")
                
                # 清理臨時文件
                try:
                    shutil.rmtree(temp_dir, ignore_errors=True)
                    shutil.rmtree(ascii_temp_dir, ignore_errors=True)
                except:
                    pass
                    
                return None
                
        except Exception as e:
            print(f"  ❌ 轉換過程中出錯: {str(e)}")
            import traceback
            traceback.print_exc()
            
            # 清理臨時文件
            try:
                shutil.rmtree(temp_dir, ignore_errors=True)
                shutil.rmtree(ascii_temp_dir, ignore_errors=True)
            except:
                pass
                
            return None
        
    except Exception as e:
        print(f"❌ 處理 PPTX 文件 {file_name} 時發生未預期的錯誤: {str(e)}")
        import traceback
        traceback.print_exc()
        return None

# 修改這個函數，確保圖片存到新的資料庫資料夾路徑
def extract_pdf_elements(file_path, file_name):
    """
    從 PDF 檔案中提取元素（文字、表格、圖片）
    file_path: PDF 檔案所在資料夾路徑
    file_name: PDF 檔案名稱
    回傳：元素列表
    """
    # 設定資料庫資料夾路徑
    DATABASE_DIR = "./database"
    output_dir = f"{DATABASE_DIR}/figures/{os.path.splitext(file_name)[0]}"
    
    # 確保輸出資料夾存在
    os.makedirs(output_dir, exist_ok=True)
    
    # 使用 partition_pdf 提取元素
    full_path = os.path.join(file_path, file_name)
    elements = partition_pdf(
        filename=full_path,
        extract_images_in_pdf=False,
        infer_table_structure=True,  
        chunking_strategy="by_title",
        max_characters=8192,
        new_after_n_chars=3800,
        combine_text_under_n_chars=2000,
        image_output_dir_path=output_dir,
        extract_image_block_output_dir=output_dir,
        strategy="hi_res",
        hi_res_model_name="yolox",
        languages=["eng","chi_tra","chi_tra_vert"],
    )
    
    return elements

def extract_pptx_elements(path, fname):
    """
    從 PowerPoint 檔案中提取元素（文字、表格、圖片）
    file_path: PPTX 檔案所在資料夾路徑
    file_name: PPTX 檔案名稱
    回傳：元素列表
    """
    # 設定資料庫資料夾路徑
    DATABASE_DIR = "./database"
    output_dir = f"{DATABASE_DIR}/figures/{os.path.splitext(fname)[0]}"
    
    # 確保輸出資料夾存在
    os.makedirs(output_dir, exist_ok=True)
    
    # 使用 partition_pptx 提取元素
    full_path = os.path.join(path, fname)
    elements = partition_pptx(
        filename=full_path,
        extract_images_in_tables=True,  
        image_output_dir_path=output_dir
    )
    
    # 清理多餘空白
    cleaned_elements = []
    for element in elements:
        # 如果元素有文字屬性，清理它
        if hasattr(element, 'text'):
            element.text = clean_extra_whitespace(element.text)
        cleaned_elements.append(element)
    
    return cleaned_elements

# 目的：方便後續針對不同型態資料進行摘要與檢索。
def categorize_elements(raw_elements):
    """
    將文件提取的元素分類為表格與文字。
    raw_elements: unstructured.documents.elements 的列表（來自PDF或PPTX）
    回傳：texts（文字列表）、tables（表格列表）
    """
    tables = []
    texts = []
    for element in raw_elements:
        element_type = str(type(element))
        if "unstructured.documents.elements.Table" in element_type:
            tables.append(str(element))
        elif "unstructured.documents.elements.CompositeElement" in element_type:
            texts.append(str(element))
        # 添加對PPTX特定元素的處理
        elif "unstructured.documents.elements.Title" in element_type:
            texts.append(str(element))
        elif "unstructured.documents.elements.NarrativeText" in element_type:
            texts.append(str(element))
    return texts, tables

if __name__ == "__main__":
    print("Wrong file execution. Please use the build_vector_db.py to build up the vector database. Then use the main.py to run the mm-RAG.")

# 添加 process_single_file 函數，從 build_vector_db.py 移植過來
def process_single_file(fname: str, fpath: str) -> dict:
    """
    處理單一檔案 (PDF 或 PPTX)，返回處理結果
    
    Args:
        fname: 檔案名稱
        fpath: 檔案所在路徑
    
    Returns:
        包含處理結果的字典:
            - texts: 提取的文字列表
            - tables: 提取的表格列表
            - page_summaries: 頁面/幻燈片摘要列表
            - page_identifiers: 頁面/幻燈片標識符列表
            - temp_files: 需要清理的臨時檔案列表
            - file_type: 檔案類型 ('pdf' 或 'pptx')
            - processing_time: 處理時間 (秒)
    """
    import time
    file_start_time = time.time()
    print(f"\n📄 處理文件: {fname}")
    
    result = {
        'texts': [],
        'tables': [],
        'page_summaries': [],
        'page_identifiers': [],
        'temp_files': [],
        'file_type': None,
        'processing_time': 0
    }
    
    try:
        if fname.lower().endswith(".pdf"):
            result['file_type'] = 'pdf'
            # 提取傳統元素
            extract_start = time.time()
            print(f"  🔄 開始提取 PDF 元素: {fname}")
            raw_elements = extract_pdf_elements(fpath, fname)
            texts, tables = categorize_elements(raw_elements)
            result['texts'] = texts
            result['tables'] = tables
            print(f"  ✓ 提取元素: {time.time() - extract_start:.2f} 秒")
            
            # 將 PDF 轉換為頁面圖片並生成摘要
            page_conversion_start = time.time()
            page_image_paths = convert_pdf_to_page_images(fpath, fname)
            print(f"  ✓ 頁面轉換: {time.time() - page_conversion_start:.2f} 秒")
            
            if page_image_paths:
                from utils.summarize import generate_pdf_page_summaries
                summary_start = time.time()
                page_summaries, page_identifiers = generate_pdf_page_summaries(fpath, fname, page_image_paths)
                result['page_summaries'] = page_summaries
                result['page_identifiers'] = page_identifiers
                print(f"  ✓ 頁面摘要生成: {time.time() - summary_start:.2f} 秒")
                
        elif fname.lower().endswith(".pptx") or fname.lower().endswith(".ppt"):
            result['file_type'] = 'pptx'
            # 提取傳統元素
            pptx_start = time.time()
            raw_elements = extract_pptx_elements(fpath, fname)
            texts, tables = categorize_elements(raw_elements)
            result['texts'] = texts
            result['tables'] = tables
            print(f"  ✓ PPTX 元素提取: {time.time() - pptx_start:.2f} 秒")
            
            # 將 PPTX 轉換為 PDF，然後處理這個 PDF
            conversion_start = time.time()
            pdf_path = convert_pptx_to_pdf(fpath, fname)
            if pdf_path:
                result['temp_files'].append(pdf_path)  # 將臨時文件加入清理列表
                
                pdf_filename = os.path.basename(pdf_path)
                pdf_dir = os.path.dirname(pdf_path)
                print(f"  ✓ PPTX 轉 PDF: {time.time() - conversion_start:.2f} 秒")
                
                # 對轉換後的 PDF 進行頁面轉換
                page_conversion_start = time.time()
                page_image_paths = convert_pdf_to_page_images(pdf_dir, pdf_filename)
                print(f"  ✓ PDF 頁面轉換: {time.time() - page_conversion_start:.2f} 秒")
                
                # 生成頁面摘要
                if page_image_paths:
                    from utils.summarize import generate_pdf_page_summaries
                    summary_start = time.time()
                    page_summaries, page_identifiers = generate_pdf_page_summaries(pdf_dir, fname, page_image_paths)
                    
                    # 修改識別符，使其反映這是 PPTX 幻燈片
                    slide_identifiers = []
                    for page_id in page_identifiers:
                        # 將 page_X 替換為 slide_X
                        slide_id = page_id.replace("_page_", "_slide_")
                        slide_identifiers.append(slide_id)
                    
                    result['page_summaries'] = page_summaries
                    result['page_identifiers'] = slide_identifiers
                    print(f"  ✓ 幻燈片摘要生成: {time.time() - summary_start:.2f} 秒")
            else:
                print(f"  ❌ PPTX 轉 PDF 失敗，無法處理幻燈片")
        
        elif fname.lower().endswith(".docx"):
            # 添加對 DOCX 文件的處理
            result['file_type'] = 'docx'
            print(f"  ⚠️ DOCX 處理功能尚未完全實現，僅提取文本")
            
            # 提取文字內容 (可以使用 unstructured 或其他庫)
            from unstructured.partition.docx import partition_docx
            
            try:
                full_path = os.path.join(fpath, fname)
                elements = partition_docx(
                    filename=full_path,
                    extract_images_in_tables=True
                )
                
                # 分類元素
                texts, tables = categorize_elements(elements)
                result['texts'] = texts
                result['tables'] = tables
                print(f"  ✓ DOCX 元素提取完成，獲取了 {len(texts)} 個文本段落和 {len(tables)} 個表格")
            except Exception as docx_err:
                print(f"  ❌ DOCX 處理出錯: {docx_err}")
                # 仍然返回部分結果
                
    except Exception as e:
        print(f"  ❌ 處理文件 {fname} 時出錯: {str(e)}")
        import traceback
        traceback.print_exc()
    
    processing_time = time.time() - file_start_time
    result['processing_time'] = processing_time
    print(f"  ✓ 文件 {fname} 處理完成: {processing_time:.2f} 秒")
    
    return result





# def process_single_file_anytype(fname: str, fpath: str) -> dict:
#     """
#     處理單一檔案，返回處理結果
    
#     Args:
#         fname: 檔案名稱
#         fpath: 檔案所在路徑
    
#     Returns:
#         包含處理結果的字典:
#             - texts: 提取的文字列表
#             - tables: 提取的表格列表
#             - page_summaries: 頁面/幻燈片摘要列表
#             - page_identifiers: 頁面/幻燈片標識符列表
#             - temp_files: 需要清理的臨時檔案列表
#             - file_type: 檔案類型 ('pdf' 或 'pptx')
#             - processing_time: 處理時間 (秒)
#     """

#     result = {
#         'texts': [],
#         'tables': [],
#         'page_summaries': [],
#         'page_identifiers': [],
#         'temp_files': [],
#         'file_type': None,
#         'processing_time': 0
#     }

#     import time
#     file_start_time = time.time()
#     print(f"\n📄 處理文件: {fname}")


#     try:
#         if fname.lower().endswith(".pdf"):
#             result['file_type'] = 'pdf'
#             table = 
        

    
#     except Exception as e:
#         print(f"  ❌ 處理文件 {fname} 時出錯: {str(e)}")
#         import traceback
#         traceback.print_exc()
    
#     processing_time = time.time() - file_start_time
#     result['processing_time'] = processing_time
#     print(f"  ✓ 文件 {fname} 處理完成: {processing_time:.2f} 秒")
    
#     return result
