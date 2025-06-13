"""PPTX related processing utilities."""

import os
import shutil
import subprocess
import tempfile
import uuid
from typing import List, Tuple

from PIL import Image
from pptx import Presentation
from unstructured.cleaners.core import clean_extra_whitespace
from unstructured.partition.pptx import partition_pptx

__all__ = [
    "convert_pptx_to_slide_images",
    "convert_pptx_to_pdf",
    "extract_pptx_elements",
]


def convert_pptx_to_slide_images(
    file_path: str, file_name: str, database_dir: str = "./database"
) -> List[Tuple[int, str]]:
    """Convert each PPTX slide to an image and return the paths."""
    try:
        output_dir = f"{database_dir}/figures/{os.path.splitext(file_name)[0]}/slides"
        os.makedirs(output_dir, exist_ok=True)
        full_path = os.path.join(file_path, file_name)
        print(f"🔄 開始處理 PPTX: {full_path}")
        if not os.path.exists(full_path):
            print(f"❌ 錯誤: PPTX 文件不存在: {full_path}")
            return []
        file_size = os.path.getsize(full_path) / (1024 * 1024)
        print(f"📄 PPTX 文件大小: {file_size:.2f} MB")
        try:
            print(f"🔄 正在載入 {file_name}...")
            presentation = Presentation(full_path)
            print(f"✅ 成功載入 {len(presentation.slides)} 張幻燈片")
        except Exception as e:
            print(f"❌ 載入 PPTX 文件時出錯: {str(e)}")
            import traceback

            traceback.print_exc()
            return []
        slide_image_paths: List[Tuple[int, str]] = []
        for i, slide in enumerate(presentation.slides):
            try:
                slide_num = i + 1
                slide_path = os.path.join(output_dir, f"slide_{slide_num}.jpg")
                slide_width, slide_height = 1920, 1080
                slide_image = Image.new(
                    "RGB", (slide_width, slide_height), color="white"
                )
                from PIL import ImageDraw, ImageFont

                draw = ImageDraw.Draw(slide_image)
                try:
                    font = ImageFont.truetype("Arial", 20)
                except IOError:
                    font = ImageFont.load_default()
                y_position = 50
                if slide.shapes.title:
                    title_text = slide.shapes.title.text
                    draw.text(
                        (50, y_position),
                        f"Title: {title_text}",
                        fill="black",
                        font=font,
                    )
                    y_position += 40
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        lines = shape.text.split("\n")
                        for line in lines:
                            if line.strip():
                                draw.text(
                                    (50, y_position), line, fill="black", font=font
                                )
                                y_position += 30
                slide_image.save(slide_path, "JPEG")
                slide_image_paths.append((slide_num, slide_path))
                print(f"  ✓ 已儲存幻燈片 {slide_num} 圖像")
            except Exception as e:
                print(f"❌ 處理幻燈片 {i+1} 時出錯: {str(e)}")
        print(
            f"✅ 已將 {file_name} 的 {len(slide_image_paths)}/{len(presentation.slides)} 張幻燈片轉換為圖片"
        )
        return slide_image_paths
    except Exception as e:
        print(f"❌ 處理 PPTX 文件 {file_name} 時發生未預期的錯誤: {str(e)}")
        import traceback

        traceback.print_exc()
        return []


def convert_pptx_to_pdf(file_path: str, file_name: str) -> str | None:
    """Convert a PPTX file to a temporary PDF and return its path."""
    try:
        full_path = os.path.join(file_path, file_name)
        if not os.path.exists(full_path):
            print(f"❌ 錯誤: PPTX 文件不存在: {full_path}")
            return None
        if not os.access(full_path, os.R_OK):
            print(f"❌ 錯誤: 沒有讀取 PPTX 文件的權限: {full_path}")
            try:
                os.chmod(full_path, 0o644)
                print("  ✓ 已嘗試修改檔案權限")
            except Exception as e:
                print(f"  ⚠️ 無法修改檔案權限: {e}")
                return None
        temp_dir = tempfile.mkdtemp()
        ascii_temp_dir = tempfile.mkdtemp()
        ascii_file_name = f"pptx_{uuid.uuid4().hex}.pptx"
        ascii_file_path = os.path.join(ascii_temp_dir, ascii_file_name)
        try:
            shutil.copy2(full_path, ascii_file_path)
            print(f"  ✓ 已創建臨時檔案副本: {ascii_file_path}")
        except Exception as e:
            print(f"  ❌ 無法創建檔案副本: {e}")
            shutil.rmtree(temp_dir, ignore_errors=True)
            shutil.rmtree(ascii_temp_dir, ignore_errors=True)
            return None
        pdf_name = f"{os.path.splitext(file_name)[0]}.pdf"
        output_pdf = os.path.join(temp_dir, pdf_name)
        ascii_output_pdf = os.path.join(
            temp_dir, f"{os.path.splitext(ascii_file_name)[0]}.pdf"
        )
        print(f"🔄 開始將 PPTX 轉換為 PDF: {ascii_file_path}")
        try:
            env = os.environ.copy()
            env["LC_ALL"] = "C"
            env["LANG"] = "C"
            conversion_methods = [
                {
                    "cmd": [
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        temp_dir,
                        ascii_file_path,
                    ],
                    "env": env,
                    "name": "標準方法",
                },
                {
                    "cmd": [
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "pdf:writer_pdf_Export",
                        "--outdir",
                        temp_dir,
                        ascii_file_path,
                    ],
                    "env": env,
                    "name": "PDF導出過濾器方法",
                },
                {
                    "cmd": [
                        "libreoffice",
                        "--headless",
                        "--infilter=impress_MS_PowerPoint_2007_XML",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        temp_dir,
                        ascii_file_path,
                    ],
                    "env": env,
                    "name": "PowerPoint過濾器方法",
                },
            ]
            success = False
            for method in conversion_methods:
                if success:
                    break
                try:
                    print(
                        f"  🔄 嘗試使用{method['name']}轉換: {' '.join(method['cmd'])}"
                    )
                    result = subprocess.run(
                        method["cmd"],
                        capture_output=True,
                        text=True,
                        check=False,
                        env=method["env"],
                        timeout=120,
                    )
                    if os.path.exists(ascii_output_pdf):
                        try:
                            shutil.move(ascii_output_pdf, output_pdf)
                            print(
                                f"  ✓ {method['name']}成功: 已將 {ascii_output_pdf} 重命名為 {output_pdf}"
                            )
                            success = True
                            break
                        except Exception as e:
                            print(f"  ⚠️ 重命名PDF文件時出錯: {e}")
                    elif os.path.exists(output_pdf):
                        print(f"  ✓ {method['name']}成功: 生成了 {output_pdf}")
                        success = True
                        break
                    else:
                        print(f"  ❌ {method['name']}失敗: 沒有生成PDF文件")
                        if result.returncode != 0:
                            print(f"  命令返回碼: {result.returncode}")
                        if result.stdout:
                            print(f"  命令輸出: {result.stdout}")
                        if result.stderr:
                            print(f"  命令錯誤: {result.stderr}")
                except subprocess.TimeoutExpired:
                    print(f"  ⚠️ {method['name']}轉換超時")
                except Exception as e:
                    print(f"  ❌ 執行{method['name']}時出錯: {e}")
            if success and os.path.exists(output_pdf):
                pdf_size = os.path.getsize(output_pdf) / 1024
                if pdf_size < 5:
                    print(
                        f"  ⚠️ 警告：生成的 PDF 文件非常小 ({pdf_size:.2f} KB)，可能轉換不完整"
                    )
                try:
                    shutil.rmtree(ascii_temp_dir, ignore_errors=True)
                except Exception:
                    pass
                return output_pdf
            print("  ❌ 所有轉換方法都失敗了")
            shutil.rmtree(temp_dir, ignore_errors=True)
            shutil.rmtree(ascii_temp_dir, ignore_errors=True)
            return None
        except Exception as e:
            print(f"  ❌ 轉換過程中出錯: {str(e)}")
            import traceback

            traceback.print_exc()
            shutil.rmtree(temp_dir, ignore_errors=True)
            shutil.rmtree(ascii_temp_dir, ignore_errors=True)
            return None
    except Exception as e:
        print(f"❌ 處理 PPTX 文件 {file_name} 時發生未預期的錯誤: {str(e)}")
        import traceback

        traceback.print_exc()
        return None


def extract_pptx_elements(path: str, fname: str, database_dir: str = "./database"):
    """Extract elements from a PPTX file."""
    output_dir = f"{database_dir}/figures/{os.path.splitext(fname)[0]}"
    os.makedirs(output_dir, exist_ok=True)
    full_path = os.path.join(path, fname)
    elements = partition_pptx(
        filename=full_path,
        extract_images_in_tables=True,
        image_output_dir_path=output_dir,
    )
    cleaned_elements = []
    for element in elements:
        if hasattr(element, "text"):
            element.text = clean_extra_whitespace(element.text)
        cleaned_elements.append(element)
    return cleaned_elements
